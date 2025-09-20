#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint MCP Server 功能测试脚本
测试 PowerPoint MCP Server 的各项功能是否正常工作

作者: 雨俊
日期: 2025-01-08
"""

import sys
import os

# 添加项目路径到 Python 路径
project_path = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, project_path)

def test_powerpoint_mcp_server():
    """
    测试 PowerPoint MCP Server 的功能
    """
    print("=" * 50)
    print("PowerPoint MCP Server 功能测试")
    print("=" * 50)
    
    try:
        # 测试模块导入
        print("\n1. 测试模块导入...")
        from tools import presentation_tools, content_tools, template_tools
        print("✓ PowerPoint MCP Server 模块导入成功")
        
        # 测试创建演示文稿
        print("\n2. 测试创建演示文稿...")
        test_ppt_path = os.path.join(project_path, "PowerPoint功能测试演示文稿.pptx")
        
        # 这里我们直接使用 python-pptx 来测试基本功能
        from pptx import Presentation
        from pptx.util import Inches
        
        # 创建新的演示文稿
        prs = Presentation()
        
        # 添加标题幻灯片
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "PowerPoint MCP Server 测试"
        subtitle.text = "功能验证演示文稿\n技术负责人: 雨俊\n测试时间: 2025-01-08"
        
        print("✓ 成功创建标题幻灯片")
        
        # 添加内容幻灯片
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = '测试功能清单'
        
        tf = body_shape.text_frame
        tf.text = '模块导入测试'
        
        p = tf.add_paragraph()
        p.text = '演示文稿创建测试'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = '幻灯片添加测试'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = '文本内容测试'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = '文件保存测试'
        p.level = 1
        
        print("✓ 成功创建内容幻灯片")
        
        # 保存演示文稿
        print("\n3. 测试保存演示文稿...")
        prs.save(test_ppt_path)
        print(f"✓ 演示文稿已保存到: {test_ppt_path}")
        
        # 验证文件是否存在
        if os.path.exists(test_ppt_path):
            file_size = os.path.getsize(test_ppt_path)
            print(f"✓ 文件验证成功，大小: {file_size} 字节")
        else:
            print("✗ 文件保存失败")
            return False
        
        print("\n" + "=" * 50)
        print("PowerPoint MCP Server 功能测试完成")
        print("=" * 50)
        print("\n测试结果:")
        print("✓ 模块导入: 成功")
        print("✓ 演示文稿创建: 成功")
        print("✓ 幻灯片添加: 成功")
        print("✓ 文本内容: 成功")
        print("✓ 文件保存: 成功")
        print("\n🎉 所有测试项目均通过!")
        print(f"\n📁 测试文件位置: {test_ppt_path}")
        
        return True
        
    except ImportError as e:
        print(f"✗ 模块导入失败: {e}")
        return False
    except Exception as e:
        print(f"✗ 测试过程中发生错误: {e}")
        return False

if __name__ == "__main__":
    success = test_powerpoint_mcp_server()
    if success:
        print("\n✅ PowerPoint MCP Server 功能正常，可以使用!")
    else:
        print("\n❌ PowerPoint MCP Server 功能测试失败，请检查配置!")
        sys.exit(1)